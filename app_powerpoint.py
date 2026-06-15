"""Étape 2 — Génération du PowerPoint.

Lit le fichier Excel nettoyé produit par `split_excel_par_slides.py`
(un onglet par slide), puis génère un fichier PowerPoint en respectant
l'ordre numérique des slides et en insérant automatiquement des slides
de transition pour les numéros sans données.
"""

import copy
import io
import re
from dataclasses import dataclass
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Chemins
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
ASSETS_DIR = BASE_DIR / "assets"
DATA_DIR = BASE_DIR / ".data"
OUTPUT_DIR = BASE_DIR / "output_ppt"

DEFAULT_TEMPLATE = ASSETS_DIR / "default.pptx"
EXAMPLE_TEMPLATE = ASSETS_DIR / "template_slide.pptx"
CLEANED_EXCEL = DATA_DIR / "Resultats_Etude_WPPmedia_DTA_clean.xlsx"
OUTPUT_BASENAME = "Powerpoint_automatise"

# ---------------------------------------------------------------------------
# Constantes de mise en page
# ---------------------------------------------------------------------------

TITLE_PLACEHOLDER_IDX = 0

LAYOUT_TITLE = "slide_title_section"
LAYOUT_GRAPH_VERTI = "slide_graph_verti"
LAYOUT_GRAPH_HORI = "slide_graph_hori"
LAYOUT_2_GRAPHS = "slide_2_graphs"
LAYOUT_3_GRAPHS = "slide_3_graphs"
LAYOUT_4_GRAPHS = "slide_4_graphs"

# ---------------------------------------------------------------------------
# Patterns de parsing des noms d'onglets (1 onglet = 1 slide, déjà normalisé)
# ---------------------------------------------------------------------------

SLIDE_SINGLE_PATTERN = re.compile(r"^[Ss]lide\s+(\d+)$")


# ---------------------------------------------------------------------------
# Modèle de données
# ---------------------------------------------------------------------------

@dataclass
class SlideData:
    """Données d'une slide issue d'un onglet Excel."""

    slide_number: int
    sheet_name: str   # nom de l'onglet (ex. "Slide 16")
    title: str        # contenu de la cellule A1 du tableau


# ---------------------------------------------------------------------------
# Lecture de l'Excel nettoyé
# ---------------------------------------------------------------------------

def parse_slide_number(sheet_name):
    """Extrait le numéro de slide depuis un nom d'onglet (ex. 'Slide 7' → 7)."""
    match = SLIDE_SINGLE_PATTERN.match(sheet_name.strip())
    if match:
        return int(match.group(1))
    raise ValueError(f"Nom d'onglet non reconnu : {sheet_name!r}")


def get_table_title_from_first_cell(worksheet):
    """Récupère le titre depuis la cellule A1 du tableau (espaces normalisés)."""
    first_cell = worksheet.cell(row=1, column=1).value
    if first_cell is None:
        return ""
    return " ".join(str(first_cell).strip().split())


def read_slides_from_excel(excel_path):
    """Lit les onglets Excel — un onglet = une slide (fichier déjà normalisé)."""
    excel_path = Path(excel_path)
    if not excel_path.exists():
        return []

    workbook = load_workbook(excel_path, read_only=True, data_only=True)
    slides_by_number = {}

    for sheet_name in workbook.sheetnames:
        worksheet = workbook[sheet_name]
        slide_number = parse_slide_number(sheet_name)
        title = get_table_title_from_first_cell(worksheet)

        if slide_number in slides_by_number:
            existing = slides_by_number[slide_number].sheet_name
            print(
                f"⚠️  Slide {slide_number} déjà définie par {existing!r}, "
                f"écrasée par {sheet_name!r}"
            )
        slides_by_number[slide_number] = SlideData(
            slide_number=slide_number,
            sheet_name=sheet_name,
            title=title or sheet_name,
        )

    workbook.close()
    return sorted(slides_by_number.values(), key=lambda s: s.slide_number)


def iter_slides_with_transitions(slides):
    """Génère la séquence complète de slides, avec transitions pour les numéros manquants.

    Pour chaque numéro entre le min et le max trouvé dans l'Excel :
    - si un onglet existe → ("content", SlideData)
    - sinon              → ("transition", numéro_manquant)
    """
    if not slides:
        return

    slides_by_number = {s.slide_number: s for s in slides}
    for n in range(min(slides_by_number), max(slides_by_number) + 1):
        if n in slides_by_number:
            yield "content", slides_by_number[n]
        else:
            yield "transition", n


# ---------------------------------------------------------------------------
# Construction du PowerPoint
# ---------------------------------------------------------------------------

def load_presentation(file_path):
    """Charge une présentation PowerPoint depuis le chemin spécifié."""
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"Template introuvable : {file_path}")
    return Presentation(str(file_path))


def get_layout_by_name(presentation, layout_name):
    """Récupère un layout par son nom ; lève une erreur claire si introuvable."""
    for layout in presentation.slide_layouts:
        if layout.name == layout_name:
            return layout
    available = ", ".join(l.name for l in presentation.slide_layouts)
    raise KeyError(f"Layout '{layout_name}' introuvable. Disponibles : {available}")


def set_slide_title(slide, title):
    """Définit le texte du placeholder de titre d'une slide s'il existe."""
    if title is None:
        return
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == TITLE_PLACEHOLDER_IDX:
            placeholder.text = title
            return


def add_slide(presentation, layout_name, title=None):
    """Ajoute une slide depuis un layout et renseigne son titre."""
    layout = get_layout_by_name(presentation, layout_name)
    slide = presentation.slides.add_slide(layout)
    set_slide_title(slide, title)
    return slide


def duplicate_slide_with_layout(presentation, original_slide, target_layout_name=LAYOUT_GRAPH_HORI):
    """Duplique une slide existante en recopiant ses formes sur un nouveau layout.

    Utilisé pour reproduire des slides d'exemple du fichier template_slide.pptx.
    """
    layout = get_layout_by_name(presentation, target_layout_name)
    duplicated_slide = presentation.slides.add_slide(layout)

    for shape in original_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_bytes = shape.image.blob
            duplicated_slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                shape.left, shape.top, shape.width, shape.height,
            )
        else:
            try:
                new_el = copy.deepcopy(shape.element)
                duplicated_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
            except Exception as exc:
                print(f"⚠️  Forme ignorée ({shape.shape_type}) : {exc}")

    return duplicated_slide


def get_unique_output_path(output_dir=OUTPUT_DIR, basename=OUTPUT_BASENAME, extension=".pptx"):
    """Renvoie un chemin de sortie disponible, avec suffixe _2, _3… si le fichier existe déjà."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    base_path = output_dir / f"{basename}{extension}"
    if not base_path.exists():
        return base_path

    counter = 2
    while True:
        candidate = output_dir / f"{basename}_{counter}{extension}"
        if not candidate.exists():
            return candidate
        counter += 1


def build_presentation(slides=None, output_file=None):
    """Construit et sauvegarde la présentation PowerPoint finale.

    - Slide 1   : titre général de l'étude (layout slide_title_section)
    - Slides 2+ : une slide par numéro dans l'ordre croissant
                  · données Excel → layout slide_graph_hori + titre depuis A1
                  · numéro manquant → slide de transition (slide_title_section)
    """
    if slides is None:
        slides = read_slides_from_excel(CLEANED_EXCEL)

    presentation = load_presentation(DEFAULT_TEMPLATE)

    add_slide(presentation, LAYOUT_TITLE, title="Résultats Étude WPPmedia DTA")

    for slide_type, slide_info in iter_slides_with_transitions(slides):
        if slide_type == "content":
            add_slide(presentation, LAYOUT_GRAPH_HORI, title=slide_info.title)
        else:
            add_slide(presentation, LAYOUT_TITLE)

    if output_file is None:
        output_file = get_unique_output_path()
    else:
        output_file = Path(output_file)
        output_file.parent.mkdir(parents=True, exist_ok=True)

    presentation.save(str(output_file))
    return output_file


# ---------------------------------------------------------------------------
# Point d'entrée
# ---------------------------------------------------------------------------

def main():
    slides = read_slides_from_excel(CLEANED_EXCEL)
    transition_count = 0

    if slides:
        print(f"📄 {len(slides)} slide(s) mappée(s) depuis {CLEANED_EXCEL.name}")
        for slide_type, slide_info in iter_slides_with_transitions(slides):
            if slide_type == "content":
                preview = slide_info.title[:80] + ("…" if len(slide_info.title) > 80 else "")
                print(f"   Slide {slide_info.slide_number:>2} ({slide_info.sheet_name}) → {preview}")
            else:
                transition_count += 1
                print(f"   Slide {slide_info:>2} → [transition]")
    else:
        print("ℹ️  Aucun Excel nettoyé trouvé : génération avec la slide de titre seule.")

    output = build_presentation(slides)
    content_count = sum(1 for t, _ in iter_slides_with_transitions(slides) if t == "content")
    total = 1 + content_count + transition_count
    print(
        f"✅ Présentation générée : {output} "
        f"({total} slide(s) : 1 titre + {content_count} contenu + {transition_count} transition)"
    )


if __name__ == "__main__":
    main()
