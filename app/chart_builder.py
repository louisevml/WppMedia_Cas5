"""
Insertion d'un graphique PowerPoint dans une diapositive.

- :const:`CHART_TYPE_MAP`           : alias ``chart_type`` du mapping → ``XL_CHART_TYPE``
- :const:`CHART_TYPES_STACKED_100`  : types qui normalisent chaque catégorie à 100 %
- :func:`get_chart_type`            : alias texte → constante python-pptx
- :func:`insert_chart_in_placeholder` : place le graphique dans le **n-ième**
  placeholder ``chart`` de la diapo (et clone depuis le layout si besoin).

Contrat côté gabarits :
    Les mises en page ``graph_verti`` / ``graph_hori`` doivent contenir un ou
    plusieurs placeholders de type « chart ». Les graphiques sont insérés
    **dans l'ordre** où ces placeholders apparaissent sur la slide ; prévoir
    donc autant de placeholders que de graphiques attendus sur une même page.
"""

from __future__ import annotations

import copy
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn


# Mapping texte (colonne chart_type) → énum python-pptx
CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "bar_clustered_vertical":     XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered_horizontal":   XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked_100_vertical":   XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar_stacked_100_horizontal": XL_CHART_TYPE.BAR_STACKED_100,
    "doughnut":                   XL_CHART_TYPE.DOUGHNUT,
    "pie":                        XL_CHART_TYPE.PIE,
}

# Types qui reçoivent une normalisation par catégorie (parts → somme 100)
CHART_TYPES_STACKED_100: frozenset[str] = frozenset({
    "bar_stacked_100_vertical",
    "bar_stacked_100_horizontal",
})


def get_chart_type(name: str) -> XL_CHART_TYPE:
    """Convertit le nom string du mapping en constante python-pptx."""
    key = name.strip().lower()
    if key not in CHART_TYPE_MAP:
        allowed = ", ".join(sorted(CHART_TYPE_MAP))
        raise ValueError(f"chart_type « {name} » inconnu. Valeurs autorisées : {allowed}")
    return CHART_TYPE_MAP[key]


# ═════════════════════════════════════════════════════════════════════════════
# Gestion des placeholders de type "chart"
# ═════════════════════════════════════════════════════════════════════════════

def _chart_placeholders_on_slide(slide) -> list:
    """Liste les placeholders chart déjà présents sur la slide."""
    out = []
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        # PP_PLACEHOLDER.CHART == 8 (certains gabarits peuvent utiliser des idx variés)
        if getattr(pf, "type", None) is not None and int(pf.type) == 8:
            out.append(ph)
    return out


def _clone_chart_placeholders_from_layout(slide) -> int:
    """Copie tous les placeholders chart du layout vers la slide. Retourne le nombre cloné."""
    layout = slide.slide_layout
    layout_spTree = layout._element.find(qn("p:cSld")).find(qn("p:spTree"))
    slide_spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))

    cloned = 0
    for sp in layout_spTree.findall(qn("p:sp")):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        ph_el = nvPr.find(qn("p:ph"))
        if ph_el is not None and ph_el.get("type") == "chart":
            slide_spTree.append(copy.deepcopy(sp))
            cloned += 1
    return cloned


def insert_chart_in_placeholder(
    slide,
    chart_type: XL_CHART_TYPE,
    chart_data: CategoryChartData,
    *,
    chart_index: int = 0,
) -> Any:
    """
    Insère un graphique dans le **n-ième** placeholder chart de la slide.

    - ``chart_index=0`` : premier placeholder chart (comportement historique,
      une slide = un graphique).
    - ``chart_index=k`` : k-ième (slide à plusieurs graphiques).

    Si aucun placeholder chart n'est encore instancié sur la slide, ils sont
    clonés depuis le layout (les placeholders définis dans le master / layout
    sont préservés lors de ``add_slide`` mais ne sont pas toujours matérialisés
    sur la slide tant qu'on n'y touche pas).

    Retourne le ``graphic_frame`` (accès au chart via ``.chart``).
    """
    placeholders = _chart_placeholders_on_slide(slide)
    if not placeholders:
        _clone_chart_placeholders_from_layout(slide)
        placeholders = _chart_placeholders_on_slide(slide)

    if not placeholders:
        raise RuntimeError(
            "Aucun placeholder de type « chart » trouvé sur cette slide "
            "ni dans sa mise en page. Ajoutez un placeholder chart au layout."
        )
    if chart_index >= len(placeholders):
        raise RuntimeError(
            f"Graphique #{chart_index + 1} demandé mais le layout ne fournit que "
            f"{len(placeholders)} placeholder(s) chart. Ajoutez un placeholder "
            "supplémentaire au layout pour accueillir ce graphique."
        )

    placeholder = placeholders[chart_index]
    return placeholder.insert_chart(chart_type, chart_data)
