"""
Gestion des gabarits PowerPoint (logique « default + banque de slides » type Ikea).

Deux fichiers ``.pptx`` sont utilisés :

- ``default.pptx`` : **présentation de sortie**. Contient le thème, les
  masques et les mises en page (slide_layouts). À chaque génération, on
  repart de ce fichier **vide** et on y ajoute les diapos une à une.

- ``template_slide.pptx`` : **banque visuelle**. Contient **6** diapos
  modèles, dans cet ordre (index 0..5) — noms de mise en page alignés
  sur le master :

  +-------+----------------------+--------------------------------+
  | index | identifiant          | usage                          |
  +=======+======================+================================+
  | 0     | ``title_section``    | Titre de section               |
  +-------+----------------------+--------------------------------+
  | 1     | ``graph_verti``      | 1 graphique « vertical »       |
  +-------+----------------------+--------------------------------+
  | 2     | ``graph_hori``       | 1 graphique « horizontal »     |
  +-------+----------------------+--------------------------------+
  | 3     | ``slide_2_graphs``   | 2 graphiques (même diapo)     |
  +-------+----------------------+--------------------------------+
  | 4     | ``slide_3_graphs``   | 3 graphiques                 |
  +-------+----------------------+--------------------------------+
  | 5     | ``slide_4_graphs``   | 4 graphiques                 |
  +-------+----------------------+--------------------------------+

Les index sont des constantes dans :mod:`app.mapping` (``TITLE_SECTION``
… ``GRAPH_4``).

Fonctions principales :

- :func:`rebuild_default_pptx_from_template` : recrée ``default.pptx`` à partir
  de ``template_slide.pptx`` (même thème + masques + mises en page, **0 diapo** document).
  À lancer après avoir ajouté de nouveaux modèles (ex. 2 / 3 / 4 graphiques) dans
  le gabarit, pour que le conteneur de génération reste aligné.
- :func:`load_default_presentation`  : charge ``default.pptx`` neuf, 0 diapo
- :func:`retrieve_slide_templates`   : indexe les 6 diapos de ``template_slide.pptx``
- :func:`add_slide_from_template`    : ajoute une diapo dans la présentation
  de sortie, en dupliquant la mise en page et les formes de la diapo modèle.
"""

from __future__ import annotations

import copy
import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .mapping import SLIDE_TYPE_NAMES


def load_presentation(path: str | Path):
    """Charge une présentation PowerPoint depuis un chemin."""
    return Presentation(str(path))


def _delete_slide_at(prs, index: int) -> None:
    """Supprime la diapositive à l'index donné (0-based)."""
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst.sldId_lst[index]
    prs.part.drop_rel(sld_id.rId)
    sld_id_lst.remove(sld_id)


def strip_all_slides(prs) -> int:
    """
    Supprime toutes les diapos document de ``prs``.

    Retourne le nombre de diapos supprimées. Utile pour produire un
    ``default.pptx`` : thème + **toutes** les définitions de mises en page
    du fichier, sans le contenu des diapos témoin.
    """
    n0 = len(prs.slides)
    while len(prs.slides):
        _delete_slide_at(prs, 0)
    return n0


def rebuild_default_pptx_from_template(
    *,
    template_slide_path: str | Path,
    default_path: str | Path,
) -> None:
    """
    Recrée le fichier **default** à partir d'une copie du **template**.

    1. Ouvre ``template_slide.pptx`` (celui où tu conçois les modèles visuels
       et, dans le master, toutes les dispositions 1 / 2 / 3 / N graphiques).
    2. Supprime chaque diapositive document — le résultat est un PPTX
       « vide » mais contenant **exactement** le même thème, masques et
       liste de ``slide_layouts`` que le fichier source.
    3. Enregistre le résultat dans ``default_path`` (défaut : ``assets/templates/default.pptx``).

    À utiliser dès que tu dupliques ou enrichis le gabarit, pour que la
    génération (qui ne charge que ``default.pptx`` comme conteneur) possède
    bien les mêmes mises en page (y compris 2, 3, 4 placeholders chart, etc.).
    """
    template_slide_path = Path(template_slide_path)
    default_path = Path(default_path)
    if not template_slide_path.is_file():
        raise FileNotFoundError(f"Fichier introuvable : {template_slide_path}")
    prs = load_presentation(template_slide_path)
    removed = strip_all_slides(prs)
    default_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(default_path))
    print(f"[OK] default.pptx ecrit : {default_path} ({default_path.stat().st_size} octets)")
    print(f"     (source: {template_slide_path}, {removed} diapos document retirees)")


def load_default_presentation(default_path: str | Path):
    """
    Ouvre ``default.pptx`` et retire toute diapositive document éventuelle
    (diapo témoin, page blanche…), afin de repartir d'un conteneur **vide**
    qui conserve uniquement thème + masques + mises en page.
    """
    prs = load_presentation(default_path)
    while len(prs.slides):
        _delete_slide_at(prs, 0)
    return prs


def retrieve_slide_templates(
    template_slide_path: str | Path,
    expected_count: int = 6,
) -> tuple[Any, dict[int, Any]]:
    """
    Charge ``template_slide.pptx`` et retourne ``(presentation, templates)``
    où ``templates[i]`` est la i-ème diapositive modèle (indexation 0-based).

    Garde la **référence** à l'objet ``Presentation`` ouvert (il faut le
    conserver en vie : les slides ne sont pas des objets détachés).

    Lève une ``ValueError`` si le fichier n'a pas au moins ``expected_count``
    diapos (convention : 6, voir le tableau dans le module docstring).
    """
    prs = load_presentation(template_slide_path)
    n = len(prs.slides)
    if n < expected_count:
        raise ValueError(
            f"template_slide.pptx doit contenir au moins {expected_count} diapositives "
            f"(title, verti, hori, 2, 3, 4 graphiques). Reçu : {n}."
        )
    templates = {idx: prs.slides[idx] for idx in range(expected_count)}
    return prs, templates


def _resolve_layout_by_name(prs, source_layout) -> Any:
    """Retrouve dans ``prs`` la mise en page portant le même nom que ``source_layout``."""
    target_name = (source_layout.name or "").strip()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if (layout.name or "").strip() == target_name:
                return layout

    names_sample = [
        (layout.name or "")
        for master in prs.slide_masters
        for layout in master.slide_layouts
    ][:24]
    raise ValueError(
        f"Disposition « {target_name!r} » introuvable dans default.pptx. "
        "default.pptx et template_slide.pptx doivent partager les mêmes noms "
        "de mises en page (exporter template_slide depuis default, ou utiliser "
        f"le même Slide Master). Layouts présents dans default : {names_sample!r}"
    )


def add_slide_from_template(prs, source_slide) -> Any:
    """
    Ajoute une nouvelle diapositive à ``prs`` (présentation de sortie) :

    1. Détermine la mise en page homologue dans ``prs`` par **nom** de layout.
    2. Crée une nouvelle diapo sur ce layout (les placeholders du layout
       — chart, titre, etc. — sont conservés).
    3. Recopie les formes **non-placeholder** de la diapo source (images,
       cadres de texte décoratifs, formes…), comme dans le pattern Ikea
       ``duplicate_slide_with_layout``.

    Retourne la nouvelle diapositive.
    """
    target_layout = _resolve_layout_by_name(prs, source_slide.slide_layout)
    new_slide = prs.slides.add_slide(target_layout)

    for shape in source_slide.shapes:
        if shape.is_placeholder:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            new_slide.shapes.add_picture(
                io.BytesIO(image.blob),
                shape.left,
                shape.top,
                width=shape.width,
                height=shape.height,
            )
            continue
        try:
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        except Exception as exc:  # pragma: no cover - best-effort copy
            print(
                f"[WARN] Forme non copiée ({shape.shape_type}) : {exc}"
            )
    return new_slide


def describe_slide_type(slide_type: int) -> str:
    """Libellé lisible pour logs/erreurs (``title_section`` / ``graph_verti`` / ``graph_hori``)."""
    return SLIDE_TYPE_NAMES.get(slide_type, f"slide_type_{slide_type}")
