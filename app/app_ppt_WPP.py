"""
Orchestrateur principal — construit le PowerPoint WPPmedia.

Même pattern que ``App_IF_gabarit`` côté Ikea :

1. Ouvrir ``default.pptx``   → présentation de sortie (vide).
2. Ouvrir ``template_slide.pptx`` → banque des 6 diapos modèles (titre, 1
   graph vert/hor, 2 / 3 / 4 graphiques).
3. Parcourir la séquence ``1..N`` construite depuis l'onglet ``mapping`` :
   - position sans ligne  → diapo titre de section.
   - position avec lignes  → modèle choisi par nombre de graphiques (1: vert/hori ;
     2 / 3 / 4 : dispositions multi-graphiques).
4. Pour les diapos graphiques, pour chaque ligne du groupe :
   - parser l'onglet Excel ``slide_N``
   - appliquer ``axis_layout``, filtres, tri, réordonnancement
   - insérer le graphique dans le n-ième placeholder chart du layout
   - appliquer le style (couleurs, légende, axes…).
5. Sauvegarder la présentation finale.
"""

from __future__ import annotations

from pathlib import Path
from typing import Callable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .chart_styles import pie_doughnut_palette
from .mapping import (
    GRAPH_2,
    GRAPH_3,
    GRAPH_4,
    GRAPH_HORI,
    GRAPH_VERTI,
    MappingRow,
    SlideSpec,
    TITLE_SECTION,
    build_slide_sequence,
    expand_split_charts,
    parse_mapping,
)
from .slide_data import (
    SlideData,
    apply_axis_layout,
    build_chart_data,
    parse_slide_sheet,
    reorder_series,
    sort_categories,
    truncate_categories_top_n,
)
from .chart_builder import (
    CHART_TYPES_STACKED_100,
    get_chart_type,
    insert_chart_in_placeholder,
)
from .slide_templates import (
    add_slide_from_template,
    describe_slide_type,
    load_default_presentation,
    retrieve_slide_templates,
)


_PIE_LIKE = {"pie", "doughnut"}
_LEGEND_LABEL_COLOR = RGBColor(64, 64, 120)
_MULTI_CHART_TYPES = {GRAPH_2, GRAPH_3, GRAPH_4}


StyleCallback = Callable[[object, str, SlideData], None]


class AppWPPgabarit:
    """
    Génère le PowerPoint WPPmedia à partir d'un Excel et des 2 gabarits.

    Paramètres :
        default_ppt_path    : chemin vers ``default.pptx``
        template_ppt_path   : chemin vers ``template_slide.pptx``
        excel_path          : fichier Excel (onglet ``mapping`` + ``slide_N``)
        output_gabarit_path : fichier de sortie
        style_callback      : fonction ``(chart, chart_type_str, slide_data) -> None``
                              appliquée après insertion de chaque graphique.
    """

    def __init__(
        self,
        *,
        default_ppt_path: str | Path,
        template_ppt_path: str | Path,
        excel_path: str | Path,
        output_gabarit_path: str | Path,
        style_callback: StyleCallback | None = None,
    ) -> None:
        self.default_ppt_path = Path(default_ppt_path)
        self.template_ppt_path = Path(template_ppt_path)
        self.excel_path = Path(excel_path)
        self.output_gabarit_path = Path(output_gabarit_path)
        self.style_callback = style_callback

        # Présentation de sortie : on part toujours d'un default « propre ».
        self.prs = load_default_presentation(self.default_ppt_path)

        # Banque visuelle : 6 diapos (indices 0..5). La Presentation source
        # doit rester en vie tant que l'on référence des slides de ce fichier.
        self._template_prs, self.template_slides = retrieve_slide_templates(
            self.template_ppt_path,
            expected_count=6,
        )

    # ── Pipeline public ─────────────────────────────────────────────────

    def run(self) -> Path:
        """Exécute le pipeline complet et retourne le chemin du PPTX généré."""
        rows = parse_mapping(self.excel_path)
        rows = expand_split_charts(self.excel_path, rows)
        if not rows:
            raise ValueError(
                "Aucune ligne de mapping trouvée — rien à générer. "
                "Vérifiez l'onglet « mapping » du fichier Excel."
            )

        sequence = build_slide_sequence(rows)
        print(
            f"[INFO] Sequence a generer : {len(sequence)} slide(s) "
            f"(titres auto-inseres pour les numeros manquants)."
        )

        for spec in sequence:
            self._generate_slide(spec)

        self.output_gabarit_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(self.output_gabarit_path))
        print(f"\n[OK] PPTX sauvegarde : {self.output_gabarit_path}")
        return self.output_gabarit_path

    # ── Étapes internes ─────────────────────────────────────────────────

    def _generate_slide(self, spec: SlideSpec) -> None:
        """Crée une diapo selon le spec (titre de section ou graphiques)."""
        source_slide = self.template_slides[spec.slide_type]
        slide = add_slide_from_template(self.prs, source_slide)

        kind = describe_slide_type(spec.slide_type)
        if spec.slide_type == TITLE_SECTION:
            print(f"[SLIDE {spec.position:>2}] {kind} (titre auto)")
            return

        nb = len(spec.rows)
        print(f"[SLIDE {spec.position:>2}] {kind} ({nb} graphique{'s' if nb > 1 else ''})")

        chart_infos: list[tuple[str, SlideData]] = []
        for chart_idx, row in enumerate(spec.rows):
            info = self._insert_graph(slide, row, chart_index=chart_idx)
            if info is not None:
                chart_infos.append(info)

        if spec.slide_type in _MULTI_CHART_TYPES and chart_infos:
            all_pie_like = all(
                ct.strip().lower() in _PIE_LIKE for ct, _ in chart_infos
            )
            if all_pie_like:
                modalities = chart_infos[0][1].categories
                if modalities and len(modalities) >= 1:
                    self._add_shared_legend(slide, modalities)

    def _insert_graph(
        self, slide, row: MappingRow, *, chart_index: int
    ) -> tuple[str, SlideData] | None:
        """Transforme les données + insère un graphique dans la diapo.

        Retourne ``(chart_type, SlideData)`` utile à l'orchestrateur pour
        construire, le cas échéant, une légende partagée (multi-donuts).
        """
        xl_chart_type = get_chart_type(row.chart_type)
        ct_key = row.chart_type.strip().lower()
        stacked_100 = ct_key in CHART_TYPES_STACKED_100

        data = parse_slide_sheet(
            self.excel_path, row.no_slide, column_pct=row.column_pct
        )
        data = apply_axis_layout(
            data,
            axis_layout=row.axis_layout,
            filter_brand=row.filter_brand,
            filter_modalite=row.filter_modalite,
            filter_segment=row.filter_segment,
        )
        if row.series_order:
            data = reorder_series(data, row.series_order)
        data = sort_categories(
            data,
            sort_order=row.sort_order,
            sort_by=row.sort_by,
            chart_type=row.chart_type,
        )
        data = truncate_categories_top_n(data, row.top_n)

        chart_data = build_chart_data(data, normalize_stacked_100=stacked_100)
        graphic_frame = insert_chart_in_placeholder(
            slide, xl_chart_type, chart_data, chart_index=chart_index,
        )
        chart = graphic_frame.chart

        if row.titre_graphique:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = row.titre_graphique
        else:
            chart.has_title = False

        if self.style_callback is not None:
            self.style_callback(chart, row.chart_type, data)

        print(
            f"    - graph #{chart_index + 1} {row.chart_type} "
            f"({len(data.categories)} categorie{'s' if len(data.categories) > 1 else ''})"
        )
        return row.chart_type, data

    def _add_shared_legend(self, slide, modalities: list[str]) -> None:
        """Ajoute une unique légende en bas de diapo : carré coloré + modalité.

        La palette utilisée est celle appliquée aux parts par ``chart_styles``.
        """
        n = len(modalities)
        palette = pie_doughnut_palette(n)

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height
        margin_x = Emu(int(slide_w * 0.04))
        legend_w = slide_w - 2 * margin_x
        legend_h = Emu(Pt(14).emu * 2)
        legend_y = Emu(int(slide_h) - int(legend_h) - Emu(int(slide_h * 0.02)))

        square_side = Pt(10)
        gap_square_text = Pt(4)
        gap_between_items = Pt(18)
        approx_char_w = Pt(5).emu
        item_widths: list[int] = []
        for name in modalities:
            text_w = int(approx_char_w * max(4, len(name) + 2))
            item_widths.append(int(square_side.emu) + int(gap_square_text.emu) + text_w)

        total_w = sum(item_widths) + int(gap_between_items.emu) * max(0, n - 1)
        start_x = int(margin_x) + max(0, (int(legend_w) - total_w) // 2)

        baseline_y = int(legend_y) + (int(legend_h) - int(square_side.emu)) // 2

        current_x = start_x
        for i, name in enumerate(modalities):
            color = palette[i] if i < len(palette) else RGBColor(128, 128, 128)
            sq = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                current_x, baseline_y,
                int(square_side.emu), int(square_side.emu),
            )
            sq.line.fill.background()
            sq.fill.solid()
            sq.fill.fore_color.rgb = color

            text_x = current_x + int(square_side.emu) + int(gap_square_text.emu)
            text_w = item_widths[i] - int(square_side.emu) - int(gap_square_text.emu)
            tb = slide.shapes.add_textbox(
                text_x,
                int(legend_y),
                text_w,
                int(legend_h),
            )
            tf = tb.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = name
            run.font.size = Pt(9)
            run.font.bold = False
            run.font.color.rgb = _LEGEND_LABEL_COLOR

            current_x += item_widths[i] + int(gap_between_items.emu)
